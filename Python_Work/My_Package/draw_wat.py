# !/usr/bin/python
# -*- coding: utf-8 -*-
# @FileName:draw_wat.py
# @Time:2020/10/11 11:36
# @Author:Jason_Yin

import re
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import matplotlib.pyplot as plt

# ---- pd设置 ----
pd.set_option('display.max_columns', None)   # 显示不省略行
pd.set_option('display.max_rows', None)      # 显示不省略列
pd.set_option('display.width', None)         # 显示不换行
# ---- plt设置 ----
plt.style.use('ggplot')    # R语言风格
# ---- regex设置 ----
regex = re.compile("\\[.*\\]")


class WatData(object):

    def __init__(self, path, datatype):
        """读取路径下的wat数据，并返回数据内容，类型，保存路径等信息"""

        data_name = "256MLPDDR2_WAT_Tracking_Table_20210420.xlsm"
        data_path = os.path.join(path, data_name)

        if datatype == "PSMC":

            spec = "psmc_spec"
            rawdata = "psmc_raw_data"
            # ---- 读取excel表中的rawdata ----
            rawdata_df = pd.read_excel(data_path, sheet_name=rawdata, header=0, index_col=0)

            rawdata_df["Test_Time"] = rawdata_df["Test_Time"].apply(lambda x: datetime.strptime(str(x), "%Y-%m-%d %H:%M:%S"))  # 时间转换

            self.save_path = os.path.join(path, 'wat_image', 'psmc')
            self.wat_type = 'DRAM'
            self.rawdata_df = rawdata_df      # 数据
            self.spec_df = pd.read_excel(data_path, sheet_name=spec, header=0, index_col=0)    # 规格

    def DrawTrend(self):
        """按时间顺序对wat的数据画trend"""

        wat_items = self.spec_df.index.values  # 获取测试项目
        for wat_item in wat_items:

            # ---- 定义变量 ----
            itemname = re.sub(regex, '', wat_item)    # 将测试项目去除单位

            # ---- 创建图层 ----
            fig = plt.figure(figsize=(8, 4))
            ax = fig.add_subplot(1, 1, 1)
            plt.subplots_adjust(left=0.1, top=0.9, bottom=0.3, right=0.9)  # 调整图片位置

            # ---- 按测试项提取参数 ----

            value = self.rawdata_df[itemname].tolist()    # Y值
            lotid = self.rawdata_df['LOT_Wafer'].tolist()     # X值
            Spec_Low = self.spec_df.loc[wat_item][0]
            Spec_Target = self.spec_df.loc[wat_item][1]
            Spec_High = self.spec_df.loc[wat_item][2]

            # ---- Y坐标轴范围 ----
            image_range = abs(Spec_Low - Spec_High)
            DrawLOL = Spec_Low - image_range*0.1
            DrawHIL = Spec_High + image_range*0.1

            # ---- 绘图 ----
            ax.plot(lotid, value, marker='o', linestyle='-', label=wat_item, linewidth=2)
            ax.axhline(y=Spec_Low, color='#FF0000', linestyle='--')
            ax.axhline(y=Spec_Target, color='#000000', linestyle='--')
            ax.axhline(y=Spec_High, color='#FF0000', linestyle='--')

            # ---- 坐标轴设置 ----

            ax.set_title("XiaoMan %s WAT (%s) Trend" % (self.wat_type, itemname))
            ax.set_xlabel('Lot_ID', fontsize=10)    # x标签设置
            ax.set_ylabel(wat_item, fontsize=10)        # y标签设置
            ax.set_ylim(DrawLOL, DrawHIL)           # y轴最大，最小值设置
            for tick in ax.get_xticklabels():       # 刻度设置
                tick.set_rotation(270)
                tick.set_fontsize(8)

            # ---- 主刻度设置 ----
            ax.yaxis.set_major_locator(plt.MaxNLocator(20))
            # ax.xaxis.set_major_locator(plt.MaxNLocator(10))

            fig.savefig(r'%s\[%s].jpg' % (self.save_path, itemname), bbox_inches='tight')
            plt.close()


def PasteToPPT(path, datatype):
    """将图片贴入ppt中"""
    file_name = "DongHu_SEDS2_WAT_Summary.pptx"
    file_path = os.path.join(path, file_name)

    #  ---- 插入图片坐标 ----
    left = [Inches(0.4), Inches(5.0), Inches(0.4), Inches(5.0)]
    top = [Inches(0.8), Inches(0.8), Inches(3.0), Inches(3.0)]
    width = [Inches(4.5), Inches(4.5), Inches(4.5), Inches(4.5)]
    height = [Inches(2), Inches(2), Inches(2), Inches(2)]

    prs = Presentation(file_path)  # 实例化

    # ---- 按工厂名称对data进行分类 ----
    if datatype == "XMC":
        image_path = os.path.join(path, r"wat_image\xmc")
        title = 'DongHu SEDS2 Logic WAT Summary'
    if datatype == "PSMC":
        image_path = os.path.join(path, r"wat_image\psmc")
        title = 'DongHu SEDS2 DRAM WAT Summary'

    image_file = dir_folder(image_path)

    for i in range(len(image_file)):
        if i % 4 == 0:
            j = 0
            slide = prs.slides.add_slide(prs.slide_layouts[2])  # 用空页布局创建一页幻灯片
            title_shape = slide.shapes.title
            title_shape.text = title
        # ---- 插入图片 ----
        slide.shapes.add_picture(image_file[i], left[j], top[j], width[j], height[j])
        j = j + 1
    prs.save(file_path)


def dir_folder(file_path):
    file_paths = []
    for root, dirs, files in os.walk(file_path):
        # root 表示当前正在访问的文件夹路径
        # dirs 表示该文件夹下的子目录名list
        # files 表示该文件夹下的文件list
        for file in files:
            if file == "Thumbs.db":   # 省略图片缩略库文件
                continue
            else:
                file_paths.append(os.path.join(root, file))
    return file_paths


if __name__ == '__main__':
    path = r"C:\Users\yinpeng\Desktop\wat_report"
    # watdata = WatData(path=path, datatype="XMC")
    # # watdata.DrawTrend()
    # PasteToPPT(path=path, datatype="XMC")
    watdata = WatData(path=path, datatype="PSMC")
    watdata.DrawTrend()
    PasteToPPT(path=path, datatype="PSMC")




